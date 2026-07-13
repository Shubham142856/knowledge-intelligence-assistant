"""
src/rag/ingestion.py — DocumentIngestor

Handles all document parsing (PDF, DOCX, PPTX, CSV, TXT), chunking into
500-token segments with 50-token overlap, and dense embedding.

Embedding backend is configurable via EMBEDDER_BACKEND env var:
  - "local"      (default) : sentence-transformers/all-MiniLM-L6-v2  → 384-dim
  - "openrouter"           : nvidia/llama-nemotron-embed-vl-1b-v2    → 2048-dim
                             (free tier via OpenRouter API)
"""

import os
import json
import logging
import time
from pathlib import Path
from typing import Union

import numpy as np
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

log = logging.getLogger("vyor_ai.ingestion")

# ── Embedding backend selection ────────────────────────────────────────────────

EMBEDDER_BACKEND = os.getenv("EMBEDDER_BACKEND", "local").lower().strip()

# Dimension map (authoritative — used by Qdrant and Titans LTM)
EMBEDDING_DIMS = {
    "local":       384,
    "openrouter":  2048,
}

EMBEDDING_DIM: int = EMBEDDING_DIMS.get(EMBEDDER_BACKEND, 384)

# ── Singleton embedder — loaded once per process, never per-request ───────────
_EMBEDDER = None


# ── OpenRouter Embedder ────────────────────────────────────────────────────────

class OpenRouterEmbedder:
    """
    Wraps OpenRouter's embeddings API endpoint.

    Provides an encode() interface matching sentence-transformers so it can
    be used as a drop-in replacement anywhere get_embedder() is called.

    Model : nvidia/llama-nemotron-embed-vl-1b-v2:free
    Dim   : 2048
    Cost  : Free tier (0 tokens cost)
    """

    DIM        = 2048
    BATCH_SIZE = 64   # max texts per API call (stay within token limits)

    def __init__(self) -> None:
        from dotenv import load_dotenv
        load_dotenv(override=True)

        self.api_key  = os.getenv("OPENROUTER_API_KEY", "")
        self.model    = os.getenv(
            "OPENROUTER_EMBED_MODEL",
            "nvidia/llama-nemotron-embed-vl-1b-v2:free",
        )
        self.site_url = os.getenv("OPENROUTER_SITE_URL", "http://localhost:8000")
        self.site_name = os.getenv("OPENROUTER_SITE_NAME", "VYOR-AI")
        self.url      = "https://openrouter.ai/api/v1/embeddings"

        if not self.api_key:
            raise ValueError(
                "OPENROUTER_API_KEY is not set. "
                "Add it to your .env file or set it as an environment variable."
            )

        import requests  # imported lazily so local mode doesn't require requests
        self._requests = requests
        log.info(
            f"OpenRouterEmbedder initialised — model={self.model}, dim={self.DIM}"
        )

    def _call_api(self, texts: list[str], retries: int = 3) -> list[list[float]]:
        """
        Call the OpenRouter embeddings endpoint.
        Supports text-only input. For multimodal, use content-array format.
        Retries up to `retries` times on transient HTTP errors.
        """
        headers = {
            "Authorization":    f"Bearer {self.api_key}",
            "Content-Type":     "application/json",
            "HTTP-Referer":     self.site_url,
            "X-Title":          self.site_name,
        }
        payload = {
            "model":           self.model,
            "input":           texts,
            "encoding_format": "float",
        }

        for attempt in range(retries):
            try:
                resp = self._requests.post(
                    self.url,
                    headers=headers,
                    data=json.dumps(payload),
                    timeout=60,
                )
                if resp.status_code == 200:
                    data = resp.json()["data"]
                    # Sort by index to preserve input order
                    data.sort(key=lambda x: x["index"])
                    return [item["embedding"] for item in data]
                elif resp.status_code == 429:
                    wait = 2 ** attempt
                    log.warning(f"OpenRouter rate limit (429). Retrying in {wait}s...")
                    time.sleep(wait)
                else:
                    log.error(
                        f"OpenRouter API error {resp.status_code}: {resp.text[:200]}"
                    )
                    resp.raise_for_status()
            except self._requests.exceptions.Timeout:
                log.warning(f"OpenRouter timeout on attempt {attempt+1}/{retries}")
                if attempt == retries - 1:
                    raise

        raise RuntimeError(f"OpenRouter embedding failed after {retries} attempts.")

    def encode(
        self,
        sentences: Union[str, list[str]],
        show_progress_bar: bool = False,
        **kwargs,
    ) -> np.ndarray:
        """
        Encode text(s) to dense embeddings.

        Matches the sentence-transformers SentenceTransformer.encode() interface.

        Args:
            sentences: A single string or list of strings.

        Returns:
            np.ndarray of shape (N, 2048) for a list, or (2048,) for a single string.
        """
        single = isinstance(sentences, str)
        if single:
            sentences = [sentences]

        all_embeddings: list[list[float]] = []

        # Batch calls to stay within rate/token limits
        for i in range(0, len(sentences), self.BATCH_SIZE):
            batch = sentences[i : i + self.BATCH_SIZE]
            if show_progress_bar:
                log.info(
                    f"  Embedding batch {i // self.BATCH_SIZE + 1}"
                    f"/{(len(sentences) + self.BATCH_SIZE - 1) // self.BATCH_SIZE}"
                )
            embeddings = self._call_api(batch)
            all_embeddings.extend(embeddings)

        arr = np.array(all_embeddings, dtype=np.float32)
        return arr[0] if single else arr

    def embed_image_text(
        self,
        text: str,
        image_url: str,
    ) -> np.ndarray:
        """
        Multimodal embedding: encode a text+image pair using the content-array format.
        Returns np.ndarray of shape (2048,).

        This is unique to the Nemotron VL model — all-MiniLM cannot do this.
        """
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type":  "application/json",
            "HTTP-Referer":  self.site_url,
            "X-Title":       self.site_name,
        }
        payload = {
            "model":           self.model,
            "input": [
                {
                    "content": [
                        {"type": "text",      "text": text},
                        {"type": "image_url", "image_url": {"url": image_url}},
                    ]
                }
            ],
            "encoding_format": "float",
        }
        resp = self._requests.post(
            self.url,
            headers=headers,
            data=json.dumps(payload),
            timeout=60,
        )
        resp.raise_for_status()
        return np.array(resp.json()["data"][0]["embedding"], dtype=np.float32)


# ── get_embedder() — singleton factory ────────────────────────────────────────

def get_embedder():
    """
    Return the module-level singleton embedder.

    Backend selection:
      EMBEDDER_BACKEND=local       → SentenceTransformer("all-MiniLM-L6-v2"),  384-dim
      EMBEDDER_BACKEND=openrouter  → OpenRouterEmbedder (nvidia/nemotron), 2048-dim
    """
    global _EMBEDDER
    if _EMBEDDER is None:
        if EMBEDDER_BACKEND == "openrouter":
            log.info("Loading OpenRouter embedder (nvidia/llama-nemotron-embed-vl-1b-v2)...")
            _EMBEDDER = OpenRouterEmbedder()
            log.info(f"OpenRouter embedder ready. dim={EMBEDDING_DIM}")
        else:
            from sentence_transformers import SentenceTransformer
            log.info("Loading sentence-transformer model all-MiniLM-L6-v2 ...")
            _EMBEDDER = SentenceTransformer("all-MiniLM-L6-v2")
            log.info("Local embedder ready.")
    return _EMBEDDER


# ── DocumentIngestor ──────────────────────────────────────────────────────────

class DocumentIngestor:
    """
    Parse → chunk → embed any supported document type.

    Supported extensions : pdf, docx, pptx, csv, txt, md
    Chunk size           : 500 tokens
    Chunk overlap        : 50 tokens
    Embedding dim        : 384  (local) | 2048 (openrouter)
    """

    CHUNK_SIZE    = 500
    CHUNK_OVERLAP = 50
    SUPPORTED     = {"pdf", "docx", "pptx", "csv", "txt", "md"}

    def __init__(self) -> None:
        self.embedder = get_embedder()
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.CHUNK_SIZE,
            chunk_overlap=self.CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

    # ── Parsers ───────────────────────────────────────────────────────────────

    def parse_pdf(self, path: str) -> str:
        """Extract text from a PDF. Raises ValueError if it appears scanned."""
        text = ""
        with open(path, "rb") as fh:
            reader = PyPDF2.PdfReader(fh)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        if len(text.strip()) < 50:
            raise ValueError(
                f"PDF appears to be scanned/image-only or empty: {path}"
            )
        return text

    def parse_docx(self, path: str) -> str:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def parse_pptx(self, path: str) -> str:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text

    def parse_csv(self, path: str) -> str:
        df = pd.read_csv(path)
        return df.to_string(index=False)

    def parse_txt(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            return fh.read()

    # ── Main entry point ──────────────────────────────────────────────────────

    def ingest(self, path: str) -> tuple[list[str], np.ndarray]:
        """
        Parse → chunk → embed a file.

        Args:
            path: Absolute path to the file.

        Returns:
            (chunks, embeddings)
            chunks:     list[str]   — Text chunks, each ≤ 500 tokens.
            embeddings: np.ndarray  — Shape (N, EMBEDDING_DIM).

        Raises:
            ValueError: Unsupported extension or scanned PDF.
        """
        ext = Path(path).suffix.lstrip(".").lower()
        if ext not in self.SUPPORTED:
            raise ValueError(
                f"Unsupported file type '.{ext}'. Allowed: {self.SUPPORTED}"
            )

        parsers = {
            "pdf":  self.parse_pdf,
            "docx": self.parse_docx,
            "pptx": self.parse_pptx,
            "csv":  self.parse_csv,
            "txt":  self.parse_txt,
            "md":   self.parse_txt,
        }

        log.info(f"Parsing {Path(path).name} [backend={EMBEDDER_BACKEND}] ...")
        raw_text = parsers[ext](path)
        log.info(f"  raw chars: {len(raw_text):,}")

        chunks = self.splitter.split_text(raw_text)
        log.info(f"  chunks: {len(chunks)}")

        embeddings = self.embedder.encode(chunks, show_progress_bar=True)
        log.info(f"  embeddings: {embeddings.shape}  (dim={embeddings.shape[-1]})")

        return chunks, embeddings
